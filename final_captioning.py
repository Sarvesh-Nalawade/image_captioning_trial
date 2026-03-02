import json
import os
import sys
from typing import Any, List
import requests

from dotenv import load_dotenv
load_dotenv()


OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')  # Set your OpenAI API key as an environment variable
OPENAI_API_URL = 'https://api.openai.com/v1/chat/completions'

def generate_detailed_image_summary(all_info: dict, retry_on_blank: bool = True) -> str:
    """
    Use OpenAI GPT-4.1 to generate a detailed image summary from all available info.
    If the API fails or returns blank, generate a fallback summary from all fields.
    """
    if not OPENAI_API_KEY:
        print("Warning: OPENAI_API_KEY not set. Skipping detailed image summary.")
        return fallback_image_summary(all_info)

    prompt = (
        "You are an expert at generating detailed, information-rich image descriptions for retrieval. "
        "Given the following information about an image, write a comprehensive, keyword-rich summary that "
        "combines all the details, so that the image can be found using any related query. "
        "Include all relevant facts, context, and text.\n\n"
        f"Page/slide text: {all_info.get('page_text') or all_info.get('slide_text', '')}\n"
        f"Text summary: {all_info.get('summary', '')}\n"
        f"BLIP caption: {all_info.get('blip_caption', '')}\n"
        f"OCR text: {'; '.join(all_info.get('ocr_text', []))}\n"
        "\nDetailed image summary:"
    )

    headers = {
        'Authorization': f'Bearer {OPENAI_API_KEY}',
        'Content-Type': 'application/json',
    }
    data = {
        "model": "gpt-4-1106-preview",  # GPT-4.1 API name
        "messages": [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt},
        ],
        "max_tokens": 256,
        "temperature": 0.4,
    }
    try:
        response = requests.post(OPENAI_API_URL, headers=headers, json=data, timeout=30)
        response.raise_for_status()
        result = response.json()
        summary = result['choices'][0]['message']['content'].strip()
        if not summary and retry_on_blank:
            print("OpenAI returned blank summary, retrying...")
            return generate_detailed_image_summary(all_info, retry_on_blank=False)
        if not summary:
            return fallback_image_summary(all_info)
        return summary
    except Exception as error:
        print(f"OpenAI API error: {error}")
        return fallback_image_summary(all_info)


def fallback_image_summary(all_info: dict) -> str:
    """
    Fallback: combine all available fields into a single summary string.
    """
    fields = []
    if all_info.get('page_text') or all_info.get('slide_text'):
        fields.append(f"Page/slide text: {all_info.get('page_text') or all_info.get('slide_text', '')}")
    if all_info.get('summary'):
        fields.append(f"Text summary: {all_info.get('summary', '')}")
    if all_info.get('blip_caption'):
        fields.append(f"BLIP caption: {all_info.get('blip_caption', '')}")
    if all_info.get('ocr_text'):
        fields.append(f"OCR text: {'; '.join(all_info.get('ocr_text', []))}")
    if not fields:
        return "No information available for this image."
    return " | ".join(fields)

from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration, pipeline
from paddleocr import PaddleOCR
from pptx import Presentation

INPUT_PATH = 'sample_files/NIV.pptx'
OUTPUT_JSON_PATH = 'sample_files/final_captioning_output.json'
PPTX_IMAGE_OUTPUT_DIR = 'sample_files/extracted_images_from_pptx'


def _extract_items_from_pptx(pptx_path: str, image_output_dir: str) -> List[dict]:
    presentation = Presentation(pptx_path)
    os.makedirs(image_output_dir, exist_ok=True)

    items: List[dict] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            text = getattr(shape, 'text', '')
            if isinstance(text, str):
                text = text.strip()
                if text:
                    slide_text_parts.append(text)
        slide_text = ' '.join(slide_text_parts)

        image_count = 0
        for shape in slide.shapes:
            if not hasattr(shape, 'image'):
                continue

            image_count += 1
            image = shape.image
            ext = image.ext.lower() if image.ext else 'png'
            image_name = f'slide{slide_index}_img{image_count}.{ext}'
            image_path = os.path.join(image_output_dir, image_name)

            with open(image_path, 'wb') as file:
                file.write(image.blob)

            items.append(
                {
                    'slide': slide_index,
                    'image_path': image_path,
                    'slide_text': slide_text,
                    'source_pptx': pptx_path,
                }
            )

    return items


def load_json(path: str) -> List[dict]:
    with open(path, 'r', encoding='utf-8') as file:
        return json.load(file)


def save_json(path: str, data: List[dict]) -> None:
    with open(path, 'w', encoding='utf-8') as file:
        json.dump(data, file, indent=2, ensure_ascii=False)


def load_input_items(path: str) -> List[dict]:
    if not os.path.exists(path):
        raise FileNotFoundError(f"Input not found: {path}")

    _, ext = os.path.splitext(path.lower())
    image_exts = {'.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.webp'}

    if ext == '.pptx':
        items = _extract_items_from_pptx(path, PPTX_IMAGE_OUTPUT_DIR)
        if not items:
            raise ValueError(f'No images found in PPTX: {path}')
        return items

    if ext in image_exts:
        return [{'image_path': path, 'page_text': ''}]

    try:
        data = load_json(path)
    except UnicodeDecodeError as error:
        raise ValueError(
            f"Input file is not valid UTF-8 JSON: {path}. "
            "Use a JSON mapping file or an image path (.png/.jpg)."
        ) from error
    except json.JSONDecodeError as error:
        raise ValueError(f"Input file is not valid JSON: {path}") from error

    if not isinstance(data, list):
        raise ValueError('Input JSON must be a list of items.')

    return data


def summarize_text(text: str, summarizer_pipe: Any) -> str:
    if not isinstance(text, str) or not text.strip():
        return ''

    source = text.strip()

    try:
        if len(source.split()) < 30:
            return source

        result = summarizer_pipe(
            source,
            max_length=80,
            min_length=20,
            do_sample=False,
            truncation=True,
        )
        return result[0]['summary_text'].strip()
    except Exception as error:
        print(f"Summarization error: {error}")
        return ''


def generate_blip_caption(image_path: str, processor: BlipProcessor, model: BlipForConditionalGeneration) -> str:
    try:
        image = Image.open(image_path).convert('RGB')
        inputs = processor(image, return_tensors='pt')
        output_tokens = model.generate(**inputs, max_new_tokens=30)
        return processor.decode(output_tokens[0], skip_special_tokens=True).strip()
    except Exception as error:
        print(f"BLIP error for {image_path}: {error}")
        return ''


def _clean_text_list(candidates: List[Any]) -> List[str]:
    seen = set()
    cleaned = []

    for item in candidates:
        if not isinstance(item, str):
            continue

        text = item.strip()
        if not text or text in seen:
            continue

        seen.add(text)
        cleaned.append(text)

    return cleaned


def extract_text_from_ocr_result(obj: Any) -> List[str]:
    text_candidates: List[str] = []

    def walk(node: Any) -> None:
        if isinstance(node, dict):
            for key, value in node.items():
                lower_key = str(key).lower()
                if lower_key in {'text', 'texts', 'rec_text', 'rec_texts', 'ocr_text'}:
                    if isinstance(value, str):
                        text_candidates.append(value)
                    elif isinstance(value, list):
                        for entry in value:
                            if isinstance(entry, str):
                                text_candidates.append(entry)
                            else:
                                walk(entry)
                    else:
                        walk(value)
                else:
                    walk(value)

        elif isinstance(node, (list, tuple)):
            if (
                len(node) >= 2
                and isinstance(node[1], (list, tuple))
                and len(node[1]) >= 1
                and isinstance(node[1][0], str)
            ):
                text_candidates.append(node[1][0])

            for item in node:
                walk(item)

    walk(obj)
    return _clean_text_list(text_candidates)


def run_ocr(image_path: str, ocr_engine: PaddleOCR) -> List[str]:
    try:
        result = ocr_engine.predict(image_path)
        texts = extract_text_from_ocr_result(result)

        if texts:
            return texts

        fallback_result = ocr_engine.ocr(image_path, cls=False)
        return extract_text_from_ocr_result(fallback_result)
    except Exception as error:
        print(f"OCR error for {image_path}: {error}")
        return []



def process_items(data: List[dict]) -> List[dict]:
    summarizer_pipe = pipeline('summarization', model='t5-small')
    blip_processor = BlipProcessor.from_pretrained('Salesforce/blip-image-captioning-base')
    blip_model = BlipForConditionalGeneration.from_pretrained('Salesforce/blip-image-captioning-base')
    ocr_engine = PaddleOCR(lang='en')

    final_items = []

    for item in data:
        image_path = item.get('image_path', '')
        page_text = item.get('page_text') or item.get('slide_text', '')

        summary = summarize_text(page_text, summarizer_pipe)

        blip_caption = ''
        ocr_text = []

        if image_path and os.path.exists(image_path):
            blip_caption = generate_blip_caption(image_path, blip_processor, blip_model)
            ocr_text = run_ocr(image_path, ocr_engine)
        else:
            print(f"Skipping missing image: {image_path}")

        merged = dict(item)
        merged['summary'] = summary
        merged['blip_caption'] = blip_caption
        merged['ocr_text'] = ocr_text

        # Generate detailed image summary using OpenAI GPT-4.1
        summary = generate_detailed_image_summary(merged)
        # Ensure no blank summary
        if not summary:
            summary = fallback_image_summary(merged)
        merged['new_image_summary'] = summary

        final_items.append(merged)

    return final_items


def main() -> None:
    input_path = sys.argv[1] if len(sys.argv) > 1 else INPUT_PATH
    data = load_input_items(input_path)

    final_output = process_items(data)
    save_json(OUTPUT_JSON_PATH, final_output)
    print(f"Processed input: {input_path}")
    print(f"Final captioning JSON saved to {OUTPUT_JSON_PATH}")


if __name__ == '__main__':
    main()
